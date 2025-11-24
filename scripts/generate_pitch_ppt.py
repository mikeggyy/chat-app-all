# -*- coding: utf-8 -*-
"""
AI 聊天應用商業簡報 PPT 生成腳本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 顏色定義
PRIMARY_COLOR = RGBColor(59, 130, 246)    # 藍色
SECONDARY_COLOR = RGBColor(16, 185, 129)  # 綠色
ACCENT_COLOR = RGBColor(139, 92, 246)     # 紫色
DARK_COLOR = RGBColor(31, 41, 55)         # 深灰
LIGHT_COLOR = RGBColor(249, 250, 251)     # 淺灰
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: 封面 =====
    slide = add_slide(prs)
    add_gradient_background(slide, PRIMARY_COLOR)

    add_text_box(slide, "AI 虛擬陪伴聊天平台",
                 Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2),
                 font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "讓每個人都能擁有專屬的 AI 夥伴",
                 Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8),
                 font_size=28, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "商業合作簡報",
                 Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5),
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== Slide 2: 一句話定位 =====
    slide = add_slide(prs)
    add_title(slide, "產品定位")

    add_text_box(slide, "「結合 AI 對話、語音互動、圖片生成的\n虛擬角色陪伴平台」",
                 Inches(0.5), Inches(1.8), Inches(12.333), Inches(1.5),
                 font_size=32, bold=True, color=PRIMARY_COLOR, align=PP_ALIGN.CENTER)

    # 核心價值
    values = [
        ("情感陪伴", "24/7 隨時對話，不會疲倦"),
        ("個性化", "可自創專屬 AI 角色"),
        ("多感官", "文字 + 語音 + 圖片互動"),
    ]

    for i, (title, desc) in enumerate(values):
        x = Inches(1 + i * 4)
        add_rounded_rect(slide, x, Inches(3.8), Inches(3.5), Inches(2),
                        fill_color=LIGHT_COLOR)
        add_text_box(slide, f"✅ {title}", x + Inches(0.2), Inches(4), Inches(3.1), Inches(0.6),
                    font_size=22, bold=True, color=DARK_COLOR)
        add_text_box(slide, desc, x + Inches(0.2), Inches(4.6), Inches(3.1), Inches(1),
                    font_size=16, color=DARK_COLOR)

    # ===== Slide 3: 市場機會 =====
    slide = add_slide(prs)
    add_title(slide, "市場機會")

    add_text_box(slide, "全球 AI 陪伴市場規模",
                 Inches(0.5), Inches(1.5), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=DARK_COLOR)

    market_data = [
        "• Character.AI：2000 萬+ 月活用戶",
        "• Replika：200 萬+ 付費用戶",
        "• 市場預計 2028 年達 $150 億美元",
        "• 年複合成長率 (CAGR) 約 23%",
    ]

    add_text_box(slide, "\n".join(market_data),
                 Inches(0.5), Inches(2.1), Inches(6), Inches(2.5),
                 font_size=18, color=DARK_COLOR)

    add_text_box(slide, "台灣/華語市場痛點",
                 Inches(7), Inches(1.5), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=DARK_COLOR)

    pain_points = [
        "• 缺乏在地化繁體中文產品",
        "• 國際產品對話體驗不佳",
        "• 文化差異導致角色不貼近",
        "• 支付方式不便（無台灣金流）",
    ]

    add_text_box(slide, "\n".join(pain_points),
                 Inches(7), Inches(2.1), Inches(6), Inches(2.5),
                 font_size=18, color=DARK_COLOR)

    # 目標用戶
    add_rounded_rect(slide, Inches(0.5), Inches(4.8), Inches(12.333), Inches(2.2),
                    fill_color=PRIMARY_COLOR)
    add_text_box(slide, "目標用戶：18-35 歲年輕族群",
                 Inches(0.7), Inches(5), Inches(12), Inches(0.5),
                 font_size=20, bold=True, color=WHITE)
    add_text_box(slide, "使用場景：深夜無人傾訴 | 情感空窗期 | 二次元愛好者 | 語言學習練習",
                 Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
                 font_size=18, color=WHITE)

    # ===== Slide 4: 核心功能 =====
    slide = add_slide(prs)
    add_title(slide, "核心功能")

    features = [
        ("💬", "智慧對話", "GPT-4o-mini 驅動\n上下文記憶\n個性化回覆", PRIMARY_COLOR),
        ("🎙️", "語音播放", "OpenAI TTS\n多種音色選擇\n即時語音生成", SECONDARY_COLOR),
        ("📸", "AI 照片", "Gemini 2.5 Flash\n角色自拍照\n即時生成", ACCENT_COLOR),
        ("✨", "自創角色", "多步驟創建流程\nAI 生成外觀\n自訂個性", RGBColor(236, 72, 153)),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        x = Inches(0.5 + i * 3.2)
        add_rounded_rect(slide, x, Inches(1.8), Inches(3), Inches(4.5),
                        fill_color=color)
        add_text_box(slide, icon, x, Inches(2), Inches(3), Inches(0.8),
                    font_size=48, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, title, x, Inches(2.9), Inches(3), Inches(0.6),
                    font_size=24, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, desc, x + Inches(0.2), Inches(3.6), Inches(2.6), Inches(2.5),
                    font_size=16, align=PP_ALIGN.CENTER, color=WHITE)

    # ===== Slide 5: 競品比較 =====
    slide = add_slide(prs)
    add_title(slide, "競品比較")

    # 表頭
    headers = ["功能", "我們", "Character.AI", "Replika"]
    col_widths = [Inches(3), Inches(2.5), Inches(2.5), Inches(2.5)]
    start_x = Inches(1.2)

    x = start_x
    for i, header in enumerate(headers):
        add_rounded_rect(slide, x, Inches(1.6), col_widths[i] - Inches(0.05), Inches(0.6),
                        fill_color=PRIMARY_COLOR)
        add_text_box(slide, header, x, Inches(1.65), col_widths[i], Inches(0.5),
                    font_size=16, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
        x += col_widths[i]

    # 資料行
    rows = [
        ["繁體中文優化", "✅ 原生", "❌ 翻譯", "❌ 翻譯"],
        ["多角色選擇", "✅ 無限", "✅ 有", "❌ 單一"],
        ["語音對話", "✅ 多音色", "❌ 無", "✅ 有"],
        ["AI 生成照片", "✅ 有", "❌ 無", "✅ 付費"],
        ["禮物互動系統", "✅ 20種", "❌ 無", "❌ 無"],
        ["等級成長系統", "✅ 有", "❌ 無", "✅ 有"],
    ]

    for row_idx, row in enumerate(rows):
        y = Inches(2.3 + row_idx * 0.7)
        bg_color = LIGHT_COLOR if row_idx % 2 == 0 else WHITE
        x = start_x
        for col_idx, cell in enumerate(row):
            add_rounded_rect(slide, x, y, col_widths[col_idx] - Inches(0.05), Inches(0.65),
                            fill_color=bg_color)
            color = SECONDARY_COLOR if "✅" in cell else (RGBColor(239, 68, 68) if "❌" in cell else DARK_COLOR)
            add_text_box(slide, cell, x, y + Inches(0.1), col_widths[col_idx], Inches(0.5),
                        font_size=14, align=PP_ALIGN.CENTER, color=color)
            x += col_widths[col_idx]

    # ===== Slide 6: 商業模式 =====
    slide = add_slide(prs)
    add_title(slide, "商業模式 - 三大收入支柱")

    pillars = [
        ("💎", "會員訂閱", "VIP $399/月\nVVIP $999/月", "預估佔比 60%", PRIMARY_COLOR),
        ("💰", "虛擬貨幣", "金幣充值\n$50 - $1,000", "預估佔比 25%", SECONDARY_COLOR),
        ("🎁", "功能購買", "AI 照片、影片\n禮物、解鎖券", "預估佔比 15%", ACCENT_COLOR),
    ]

    for i, (icon, title, desc, ratio, color) in enumerate(pillars):
        x = Inches(0.8 + i * 4.2)
        add_rounded_rect(slide, x, Inches(1.8), Inches(3.8), Inches(4.2),
                        fill_color=color)
        add_text_box(slide, icon, x, Inches(2), Inches(3.8), Inches(0.8),
                    font_size=48, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, title, x, Inches(2.9), Inches(3.8), Inches(0.6),
                    font_size=26, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, desc, x, Inches(3.6), Inches(3.8), Inches(1.5),
                    font_size=18, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, ratio, x, Inches(5.2), Inches(3.8), Inches(0.5),
                    font_size=16, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

    # ===== Slide 7: 會員方案 =====
    slide = add_slide(prs)
    add_title(slide, "會員方案設計")

    # VVIP
    add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(4.8),
                    fill_color=ACCENT_COLOR)
    add_text_box(slide, "👑 VVIP", Inches(0.5), Inches(1.8), Inches(4), Inches(0.6),
                font_size=28, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "NT$ 999/月", Inches(0.5), Inches(2.4), Inches(4), Inches(0.5),
                font_size=22, align=PP_ALIGN.CENTER, color=WHITE)
    vvip_features = "• 50 次對話/角色\n• 無限配對\n• GPT-4.1-mini\n• 金幣 8 折\n\n開通禮：\n2000 金幣\n30 解鎖券\n60 拍照卡"
    add_text_box(slide, vvip_features, Inches(0.7), Inches(3.1), Inches(3.6), Inches(3),
                font_size=14, color=WHITE)

    # VIP
    add_rounded_rect(slide, Inches(4.7), Inches(1.6), Inches(4), Inches(4.8),
                    fill_color=PRIMARY_COLOR)
    add_text_box(slide, "⭐ VIP", Inches(4.7), Inches(1.8), Inches(4), Inches(0.6),
                font_size=28, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "NT$ 399/月", Inches(4.7), Inches(2.4), Inches(4), Inches(0.5),
                font_size=22, align=PP_ALIGN.CENTER, color=WHITE)
    vip_features = "• 20 次對話/角色\n• 30 次配對/日\n• 無限語音\n• 金幣 9 折\n\n開通禮：\n600 金幣\n10 解鎖券\n20 拍照卡"
    add_text_box(slide, vip_features, Inches(4.9), Inches(3.1), Inches(3.6), Inches(3),
                font_size=14, color=WHITE)

    # Free
    add_rounded_rect(slide, Inches(8.9), Inches(1.6), Inches(4), Inches(4.8),
                    fill_color=RGBColor(107, 114, 128))
    add_text_box(slide, "🆓 免費會員", Inches(8.9), Inches(1.8), Inches(4), Inches(0.6),
                font_size=28, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "NT$ 0", Inches(8.9), Inches(2.4), Inches(4), Inches(0.5),
                font_size=22, align=PP_ALIGN.CENTER, color=WHITE)
    free_features = "• 10 次對話/角色\n• 5 次配對/日\n• 10 次語音/角色\n• 可看廣告解鎖\n\n轉化機制：\n達到限制時\n引導升級 VIP"
    add_text_box(slide, free_features, Inches(9.1), Inches(3.1), Inches(3.6), Inches(3),
                font_size=14, color=WHITE)

    # ===== Slide 8: 虛擬經濟 =====
    slide = add_slide(prs)
    add_title(slide, "虛擬經濟系統")

    # 金幣套餐
    add_text_box(slide, "💰 金幣充值套餐", Inches(0.5), Inches(1.5), Inches(6), Inches(0.5),
                font_size=20, bold=True, color=DARK_COLOR)

    packages = [
        ("100 金幣", "NT$ 50", ""),
        ("500 金幣", "NT$ 200", "+50 贈送"),
        ("1000 金幣", "NT$ 350", "+150 贈送 推薦"),
        ("3000 金幣", "NT$ 1,000", "+500 贈送 最超值"),
    ]

    for i, (coins, price, bonus) in enumerate(packages):
        y = Inches(2.1 + i * 0.6)
        add_text_box(slide, f"{coins}  =  {price}  {bonus}",
                    Inches(0.7), y, Inches(5.5), Inches(0.5),
                    font_size=16, color=DARK_COLOR)

    # 消費用途
    add_text_box(slide, "🎯 金幣消費用途", Inches(7), Inches(1.5), Inches(6), Inches(0.5),
                font_size=20, bold=True, color=DARK_COLOR)

    usages = [
        ("📸 AI 自拍照", "50 金幣"),
        ("🎬 AI 短影片", "200 金幣"),
        ("🔓 角色解鎖", "300 金幣"),
        ("🎁 虛擬禮物", "10-2000 金幣"),
    ]

    for i, (item, price) in enumerate(usages):
        y = Inches(2.1 + i * 0.6)
        add_text_box(slide, f"{item}：{price}",
                    Inches(7.2), y, Inches(5.5), Inches(0.5),
                    font_size=16, color=DARK_COLOR)

    # 禮物系統
    add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(2),
                    fill_color=LIGHT_COLOR)
    add_text_box(slide, "🎁 禮物系統（20 種）", Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
                font_size=18, bold=True, color=DARK_COLOR)
    add_text_box(slide, "普通 10-25 → 罕見 30-60 → 稀有 70-120 → 史詩 150-300 → 傳說 400-2000 金幣",
                Inches(0.7), Inches(5.3), Inches(12), Inches(0.5),
                font_size=16, color=DARK_COLOR)
    add_text_box(slide, "禮物可增加好感度，提升角色等級，解鎖專屬互動內容",
                Inches(0.7), Inches(5.8), Inches(12), Inches(0.5),
                font_size=14, color=RGBColor(107, 114, 128))

    # ===== Slide 9: 用戶轉化漏斗 =====
    slide = add_slide(prs)
    add_title(slide, "用戶轉化漏斗")

    # 漏斗圖
    funnel = [
        ("訪客體驗", "2 次對話限制", Inches(12), RGBColor(209, 213, 219)),
        ("免費註冊", "10 次對話 + 廣告解鎖", Inches(10), RGBColor(156, 163, 175)),
        ("VIP 轉化", "NT$ 399/月", Inches(7), PRIMARY_COLOR),
        ("VVIP 轉化", "NT$ 999/月", Inches(4), ACCENT_COLOR),
    ]

    for i, (stage, desc, width, color) in enumerate(funnel):
        x = (Inches(13.333) - width) / 2
        y = Inches(1.8 + i * 1.2)
        add_rounded_rect(slide, x, y, width, Inches(1), fill_color=color)
        text_color = WHITE if i >= 2 else DARK_COLOR
        add_text_box(slide, f"{stage}  |  {desc}", x, y + Inches(0.25), width, Inches(0.5),
                    font_size=18, bold=True, align=PP_ALIGN.CENTER, color=text_color)

    # 轉化率
    add_text_box(slide, "預期轉化率", Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
                font_size=16, bold=True, color=DARK_COLOR)
    add_text_box(slide, "註冊→活躍: 60%   |   活躍→付費: 5-10%   |   VIP→VVIP: 15-20%",
                Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
                font_size=14, color=RGBColor(107, 114, 128))

    # ===== Slide 10: 技術架構 =====
    slide = add_slide(prs)
    add_title(slide, "技術架構")

    # 架構層
    layers = [
        ("前端 (Vue 3 + Vite)", "響應式設計、PWA 支援", Inches(1.5), SECONDARY_COLOR),
        ("後端 API (Node.js + Express)", "完整測試覆蓋、速率限制", Inches(2.7), PRIMARY_COLOR),
        ("外部服務", "OpenAI GPT-4 | OpenAI TTS | Gemini 圖片 | Replicate 影片", Inches(3.9), ACCENT_COLOR),
        ("數據層 (Firebase)", "Firestore + Auth + Storage", Inches(5.1), RGBColor(251, 146, 60)),
    ]

    for label, desc, y, color in layers:
        add_rounded_rect(slide, Inches(1), y, Inches(11.333), Inches(1), fill_color=color)
        add_text_box(slide, label, Inches(1.2), y + Inches(0.15), Inches(5), Inches(0.4),
                    font_size=16, bold=True, color=WHITE)
        add_text_box(slide, desc, Inches(6), y + Inches(0.15), Inches(6), Inches(0.4),
                    font_size=14, color=WHITE)

    # ===== Slide 11: 技術亮點 =====
    slide = add_slide(prs)
    add_title(slide, "技術亮點")

    highlights = [
        ("🔒", "企業級安全", "Firebase Auth\nAPI 速率限制\n冪等性保護\n日誌脫敏"),
        ("📊", "高品質代碼", "688 個自動化測試\n31 個 API 覆蓋\n100% 測試通過"),
        ("⚡", "性能優化", "多層緩存系統\n圖片壓縮 70-85%\n虛擬滾動"),
        ("🛠️", "可擴展架構", "前後端分離\n模組化設計\n完整管理後臺"),
    ]

    for i, (icon, title, desc) in enumerate(highlights):
        x = Inches(0.5 + i * 3.2)
        add_rounded_rect(slide, x, Inches(1.8), Inches(3), Inches(4.2),
                        fill_color=LIGHT_COLOR)
        add_text_box(slide, icon, x, Inches(2), Inches(3), Inches(0.7),
                    font_size=40, align=PP_ALIGN.CENTER, color=DARK_COLOR)
        add_text_box(slide, title, x, Inches(2.7), Inches(3), Inches(0.5),
                    font_size=20, bold=True, align=PP_ALIGN.CENTER, color=DARK_COLOR)
        add_text_box(slide, desc, x + Inches(0.2), Inches(3.4), Inches(2.6), Inches(2.5),
                    font_size=14, align=PP_ALIGN.CENTER, color=DARK_COLOR)

    # ===== Slide 12: 發展規劃 =====
    slide = add_slide(prs)
    add_title(slide, "發展規劃")

    phases = [
        ("Phase 1", "穩定運營", "✅ 核心功能完成\n✅ 商業模式驗證\n✅ 技術架構成熟", "現階段", SECONDARY_COLOR),
        ("Phase 2", "用戶增長", "□ 社群推廣\n□ KOL 合作\n□ App Store 上架", "3-6 個月", PRIMARY_COLOR),
        ("Phase 3", "功能擴展", "□ 即時語音對話\n□ 多語言支援\n□ AR 虛擬角色", "6-12 個月", ACCENT_COLOR),
        ("Phase 4", "生態擴展", "□ 創作者分潤\n□ UGC 角色市場\n□ API 開放平台", "12+ 個月", RGBColor(236, 72, 153)),
    ]

    for i, (phase, title, desc, timeline, color) in enumerate(phases):
        x = Inches(0.3 + i * 3.3)
        add_rounded_rect(slide, x, Inches(1.7), Inches(3.1), Inches(4.5),
                        fill_color=color)
        add_text_box(slide, phase, x, Inches(1.85), Inches(3.1), Inches(0.4),
                    font_size=14, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, title, x, Inches(2.25), Inches(3.1), Inches(0.5),
                    font_size=20, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, desc, x + Inches(0.2), Inches(2.9), Inches(2.7), Inches(2.5),
                    font_size=13, color=WHITE)
        add_text_box(slide, timeline, x, Inches(5.5), Inches(3.1), Inches(0.4),
                    font_size=12, align=PP_ALIGN.CENTER, color=WHITE)

    # ===== Slide 13: 合作需求 =====
    slide = add_slide(prs)
    add_title(slide, "尋求合作")

    # 資金需求
    add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(2.8),
                    fill_color=PRIMARY_COLOR)
    add_text_box(slide, "💰 資金需求", Inches(0.5), Inches(1.8), Inches(4), Inches(0.5),
                font_size=20, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "• 伺服器與 AI API 成本\n• 用戶獲取與推廣\n• 團隊擴充",
                Inches(0.7), Inches(2.5), Inches(3.6), Inches(1.8),
                font_size=16, color=WHITE)

    # 資源需求
    add_rounded_rect(slide, Inches(4.7), Inches(1.6), Inches(4), Inches(2.8),
                    fill_color=SECONDARY_COLOR)
    add_text_box(slide, "🤝 資源需求", Inches(4.7), Inches(1.8), Inches(4), Inches(0.5),
                font_size=20, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "• 公司名義（金流/發票）\n• 社群運營經驗\n• App 上架與推廣",
                Inches(4.9), Inches(2.5), Inches(3.6), Inches(1.8),
                font_size=16, color=WHITE)

    # 回報
    add_rounded_rect(slide, Inches(8.9), Inches(1.6), Inches(4), Inches(2.8),
                    fill_color=ACCENT_COLOR)
    add_text_box(slide, "📈 預期效益", Inches(8.9), Inches(1.8), Inches(4), Inches(0.5),
                font_size=20, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, "• 技術成熟，可快速上線\n• 多元變現模式\n• 可擴展的商業潛力",
                Inches(9.1), Inches(2.5), Inches(3.6), Inches(1.8),
                font_size=16, color=WHITE)

    # 合作模式
    add_rounded_rect(slide, Inches(0.5), Inches(4.7), Inches(12.333), Inches(1.8),
                    fill_color=LIGHT_COLOR)
    add_text_box(slide, "合作模式選項", Inches(0.7), Inches(4.9), Inches(12), Inches(0.4),
                font_size=18, bold=True, color=DARK_COLOR)
    add_text_box(slide, "方案 A：借用公司名義（分潤模式）  |  方案 B：入股合作（股權模式）",
                Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
                font_size=16, color=DARK_COLOR)
    add_text_box(slide, "具體條款可進一步討論",
                Inches(0.7), Inches(5.9), Inches(12), Inches(0.4),
                font_size=14, color=RGBColor(107, 114, 128))

    # ===== Slide 14: 為什麼值得投資 =====
    slide = add_slide(prs)
    add_title(slide, "為什麼值得合作？")

    reasons = [
        ("1️⃣", "市場時機對", "AI 陪伴需求爆發期，華語市場尚未飽和"),
        ("2️⃣", "產品成熟度高", "核心功能完成，可立即上線運營"),
        ("3️⃣", "技術門檻高", "688 個測試、企業級架構，不易被快速複製"),
        ("4️⃣", "變現模式清晰", "會員訂閱 + 虛擬貨幣 + 功能購買，多元收入"),
        ("5️⃣", "可擴展性強", "有明確的產品路線圖和成長路徑"),
    ]

    for i, (num, title, desc) in enumerate(reasons):
        y = Inches(1.6 + i * 1.1)
        add_rounded_rect(slide, Inches(0.5), y, Inches(12.333), Inches(0.95),
                        fill_color=LIGHT_COLOR if i % 2 == 0 else WHITE)
        add_text_box(slide, num, Inches(0.7), y + Inches(0.2), Inches(0.8), Inches(0.5),
                    font_size=24, color=PRIMARY_COLOR)
        add_text_box(slide, title, Inches(1.5), y + Inches(0.15), Inches(3), Inches(0.5),
                    font_size=20, bold=True, color=DARK_COLOR)
        add_text_box(slide, desc, Inches(4.5), y + Inches(0.2), Inches(8), Inches(0.5),
                    font_size=16, color=DARK_COLOR)

    # ===== Slide 15: 結尾 =====
    slide = add_slide(prs)
    add_gradient_background(slide, PRIMARY_COLOR)

    add_text_box(slide, "感謝聆聽！",
                 Inches(0.5), Inches(2.5), Inches(12.333), Inches(1),
                 font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "期待與您合作，一起打造 AI 陪伴新體驗",
                 Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.6),
                 font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "（請填入您的聯繫方式）\n電話：\nEmail：\nLine ID：",
                 Inches(0.5), Inches(5), Inches(12.333), Inches(1.5),
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

    return prs


def add_slide(prs):
    """添加空白投影片"""
    blank_layout = prs.slide_layouts[6]  # 空白版面
    return prs.slides.add_slide(blank_layout)


def add_title(slide, title_text):
    """添加標題"""
    add_text_box(slide, title_text, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8),
                font_size=36, bold=True, color=DARK_COLOR)
    # 底線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(2), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()


def add_text_box(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_COLOR, align=PP_ALIGN.LEFT):
    """添加文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_COLOR):
    """添加圓角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # 調整圓角大小
    shape.adjustments[0] = 0.1
    return shape


def add_gradient_background(slide, color):
    """添加漸層背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # 移到最底層
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


if __name__ == "__main__":
    print("正在生成 PPT 簡報...")
    prs = create_presentation()

    output_path = os.path.join(os.path.dirname(__file__), "..", "AI_Chat_App_Pitch.pptx")
    prs.save(output_path)

    print(f"✅ PPT 已生成：{os.path.abspath(output_path)}")
    print(f"   共 {len(prs.slides)} 張投影片")
